#!/usr/bin/env python3
"""
报告生成模块
功能：PDF报告生成、PPT报告生成、Word报告生成
"""

import io
import logging
import os
from dataclasses import dataclass
from datetime import datetime
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (
        Image,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    logger.warning("ReportLab未安装，PDF生成功能不可用")

try:
    from pptx import Presentation
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx未安装，PPT生成功能不可用")


@dataclass
class ReportConfig:
    """报告配置"""
    title: str = "舆情分析报告"
    subtitle: str = ""
    author: str = "微博舆情分析系统"
    date_format: str = "%Y年%m月%d日"
    include_charts: bool = True
    include_tables: bool = True
    page_size: str = "A4"


class PDFReportGenerator:
    """PDF报告生成器"""

    def __init__(self):
        self.styles = None
        self._init_styles()

    def _init_styles(self):
        """初始化样式"""
        if not REPORTLAB_AVAILABLE:
            return

        self.styles = getSampleStyleSheet()

        try:
            font_path = self._find_chinese_font()
            if font_path:
                pdfmetrics.registerFont(TTFont('Chinese', font_path))
                self.chinese_font = 'Chinese'
            else:
                self.chinese_font = 'Helvetica'
        except Exception:
            self.chinese_font = 'Helvetica'

        self.styles.add(ParagraphStyle(
            name='ChineseTitle',
            fontName=self.chinese_font,
            fontSize=24,
            leading=30,
            alignment=1,
            spaceAfter=20
        ))

        self.styles.add(ParagraphStyle(
            name='ChineseHeading',
            fontName=self.chinese_font,
            fontSize=16,
            leading=20,
            spaceBefore=15,
            spaceAfter=10
        ))

        self.styles.add(ParagraphStyle(
            name='ChineseBody',
            fontName=self.chinese_font,
            fontSize=11,
            leading=16,
            spaceBefore=6,
            spaceAfter=6
        ))

    def _find_chinese_font(self) -> Optional[str]:
        """查找中文字体"""
        font_paths = [
            "C:/Windows/Fonts/simhei.ttf",
            "C:/Windows/Fonts/msyh.ttc",
            "C:/Windows/Fonts/simsun.ttc",
            "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        ]

        for path in font_paths:
            if os.path.exists(path):
                return path
        return None

    def generate(
        self,
        data: Dict[str, Any],
        output_path: str = None,
        config: ReportConfig = None
    ) -> Optional[str]:
        """
        生成PDF报告

        Args:
            data: 报告数据
            output_path: 输出路径
            config: 报告配置

        Returns:
            str: 生成的文件路径
        """
        if not REPORTLAB_AVAILABLE:
            logger.error("ReportLab未安装")
            return None

        config = config or ReportConfig()

        if output_path is None:
            output_path = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"

        try:
            doc = SimpleDocTemplate(
                output_path,
                pagesize=A4,
                rightMargin=2*cm,
                leftMargin=2*cm,
                topMargin=2*cm,
                bottomMargin=2*cm
            )

            story = []

            story.append(Paragraph(config.title, self.styles['ChineseTitle']))
            story.append(Paragraph(config.subtitle, self.styles['ChineseBody']))
            story.append(Paragraph(
                f"生成日期: {datetime.now().strftime(config.date_format)}",
                self.styles['ChineseBody']
            ))
            story.append(Spacer(1, 30))

            if 'summary' in data:
                story.append(Paragraph("一、数据概览", self.styles['ChineseHeading']))
                summary = data['summary']

                summary_data = [
                    ['指标', '数值'],
                    ['总文章数', str(summary.get('total_articles', 0))],
                    ['总评论数', str(summary.get('total_comments', 0))],
                    ['正面评价', str(summary.get('positive_count', 0))],
                    ['中性评价', str(summary.get('neutral_count', 0))],
                    ['负面评价', str(summary.get('negative_count', 0))],
                ]

                table = Table(summary_data, colWidths=[8*cm, 6*cm])
                table.setStyle(TableStyle([
                    ('FONTNAME', (0, 0), (-1, -1), self.chinese_font),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ]))
                story.append(table)
                story.append(Spacer(1, 20))

            if 'sentiment_analysis' in data:
                story.append(Paragraph("二、情感分析", self.styles['ChineseHeading']))
                sentiment = data['sentiment_analysis']

                for key, value in sentiment.items():
                    story.append(Paragraph(
                        f"• {key}: {value}",
                        self.styles['ChineseBody']
                    ))
                story.append(Spacer(1, 20))

            if 'hot_topics' in data:
                story.append(Paragraph("三、热门话题", self.styles['ChineseHeading']))

                topics_data = [['排名', '话题', '热度']]
                for i, topic in enumerate(data['hot_topics'][:10], 1):
                    topics_data.append([
                        str(i),
                        topic.get('name', '')[:30],
                        str(topic.get('heat', 0))
                    ])

                table = Table(topics_data, colWidths=[2*cm, 10*cm, 3*cm])
                table.setStyle(TableStyle([
                    ('FONTNAME', (0, 0), (-1, -1), self.chinese_font),
                    ('FONTSIZE', (0, 0), (-1, -1), 9),
                    ('BACKGROUND', (0, 0), (-1, 0), colors.Color(0.2, 0.4, 0.6)),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ]))
                story.append(table)
                story.append(Spacer(1, 20))

            if 'alerts' in data:
                story.append(Paragraph("四、预警记录", self.styles['ChineseHeading']))

                for alert in data['alerts'][:10]:
                    story.append(Paragraph(
                        f"• [{alert.get('level', 'info')}] {alert.get('title', '')}: {alert.get('message', '')}",
                        self.styles['ChineseBody']
                    ))

            story.append(Spacer(1, 30))
            story.append(Paragraph(
                f"报告生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
                self.styles['ChineseBody']
            ))

            doc.build(story)
            logger.info(f"PDF报告已生成: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"PDF生成失败: {e}")
            return None


class PPTReportGenerator:
    """PPT报告生成器"""

    def __init__(self):
        self.prs = None

    def generate(
        self,
        data: Dict[str, Any],
        output_path: str = None,
        config: ReportConfig = None
    ) -> Optional[str]:
        """
        生成PPT报告
        """
        if not PPTX_AVAILABLE:
            logger.error("python-pptx未安装")
            return None

        config = config or ReportConfig()

        if output_path is None:
            output_path = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

        try:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)

            self._add_title_slide(config)

            if 'summary' in data:
                self._add_summary_slide(data['summary'])

            if 'sentiment_analysis' in data:
                self._add_sentiment_slide(data['sentiment_analysis'])

            if 'hot_topics' in data:
                self._add_topics_slide(data['hot_topics'])

            if 'alerts' in data:
                self._add_alerts_slide(data['alerts'])

            self._add_end_slide()

            self.prs.save(output_path)
            logger.info(f"PPT报告已生成: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"PPT生成失败: {e}")
            return None

    def _add_title_slide(self, config: ReportConfig):
        """添加标题幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = config.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"生成日期: {datetime.now().strftime(config.date_format)}"
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER

    def _add_summary_slide(self, summary: Dict):
        """添加概览幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "数据概览"
        p.font.size = Pt(32)
        p.font.bold = True

        metrics = [
            ('总文章数', summary.get('total_articles', 0)),
            ('总评论数', summary.get('total_comments', 0)),
            ('正面评价', summary.get('positive_count', 0)),
            ('中性评价', summary.get('neutral_count', 0)),
            ('负面评价', summary.get('negative_count', 0)),
        ]

        for i, (label, value) in enumerate(metrics):
            left = Inches(0.5 + (i % 3) * 4.2)
            top = Inches(1.5 + (i // 3) * 2.5)

            box = slide.shapes.add_textbox(left, top, Inches(3.8), Inches(2))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(48)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            p2 = tf.add_paragraph()
            p2.text = label
            p2.font.size = Pt(18)
            p2.alignment = PP_ALIGN.CENTER

    def _add_sentiment_slide(self, sentiment: Dict):
        """添加情感分析幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "情感分析"
        p.font.size = Pt(32)
        p.font.bold = True

        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)
        )
        tf = content_box.text_frame

        for key, value in sentiment.items():
            p = tf.add_paragraph()
            p.text = f"• {key}: {value}"
            p.font.size = Pt(20)
            p.space_after = Pt(12)

    def _add_topics_slide(self, topics: List[Dict]):
        """添加热门话题幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "热门话题 Top 10"
        p.font.size = Pt(32)
        p.font.bold = True

        rows = min(len(topics), 10)
        cols = 2

        for i, topic in enumerate(topics[:10]):
            left = Inches(0.5 + (i % 2) * 6.4)
            top = Inches(1.3 + (i // 2) * 1.1)

            box = slide.shapes.add_textbox(left, top, Inches(6), Inches(0.9))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{i+1}. {topic.get('name', '')[:25]}"
            p.font.size = Pt(18)

            p2 = tf.add_paragraph()
            p2.text = f"热度: {topic.get('heat', 0)}"
            p2.font.size = Pt(14)
            p2.font.color.rgb = RgbColor(100, 100, 100)

    def _add_alerts_slide(self, alerts: List[Dict]):
        """添加预警幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "预警记录"
        p.font.size = Pt(32)
        p.font.bold = True

        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5)
        )
        tf = content_box.text_frame

        for alert in alerts[:8]:
            p = tf.add_paragraph()
            level = alert.get('level', 'info')
            p.text = f"[{level.upper()}] {alert.get('title', '')}"
            p.font.size = Pt(16)
            p.space_after = Pt(8)

            p2 = tf.add_paragraph()
            p2.text = f"  {alert.get('message', '')}"
            p2.font.size = Pt(14)
            p2.font.color.rgb = RgbColor(100, 100, 100)
            p2.space_after = Pt(16)

    def _add_end_slide(self):
        """添加结束幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "谢谢"
        p.font.size = Pt(48)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER


class ReportGenerator:
    """统一报告生成器"""

    def __init__(self):
        self.pdf_generator = PDFReportGenerator()
        self.ppt_generator = PPTReportGenerator()

    def generate_report(
        self,
        data: Dict[str, Any],
        format: str = 'pdf',
        output_path: str = None,
        config: ReportConfig = None
    ) -> Optional[str]:
        """
        生成报告

        Args:
            data: 报告数据
            format: 格式
            output_path: 输出路径
            config: 配置

        Returns:
            str: 生成的文件路径
        """
        if format.lower() == 'pdf':
            return self.pdf_generator.generate(data, output_path, config)
        elif format.lower() == 'ppt':
            return self.ppt_generator.generate(data, output_path, config)
        else:
            logger.error(f"不支持的格式: {format}")
            return None

    def generate_all(
        self,
        data: Dict[str, Any],
        output_dir: str = None,
        config: ReportConfig = None
    ) -> Dict[str, str]:
        """
        生成所有格式报告

        Returns:
            Dict[str, str]: 格式到文件路径的映射
        """
        results = {}

        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            pdf_path = os.path.join(output_dir, f"report_{timestamp}.pdf")
            ppt_path = os.path.join(output_dir, f"report_{timestamp}.pptx")
        else:
            pdf_path = None
            ppt_path = None

        pdf_result = self.generate_report(data, 'pdf', pdf_path, config)
        if pdf_result:
            results['pdf'] = pdf_result

        ppt_result = self.generate_report(data, 'ppt', ppt_path, config)
        if ppt_result:
            results['ppt'] = ppt_result

        return results


report_generator = ReportGenerator()
